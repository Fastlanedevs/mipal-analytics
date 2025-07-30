import io
import logging
import asyncio
import base64
from typing import Optional, Dict, Any

import magic
import pdfplumber
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pydantic import BaseModel
from pkg.log.logger import Logger


# -----------------------------------------------------------------------------
# Pydantic Model for Extraction Result
# -----------------------------------------------------------------------------
class ExtractedContent(BaseModel):
    file_id: str
    mime_type: str
    extracted_text: Optional[str] = None
    structured_data: Optional[Any] = None
    raw_content: Optional[bytes] = None


# -----------------------------------------------------------------------------
# File Content Extractor: Converts Bytes to Text
# -----------------------------------------------------------------------------
class ContentExtractor:
    def __init__(self, logger: Logger):
        # Use python-magic to detect the MIME type from bytes.
        self.mime_detector = magic.Magic(mime=True)
        self.logger = logger


    def _get_handler(self, mime_type: str):
        """
        Returns the appropriate extraction handler based on the MIME type.
        Handles both original Google MIME types and exported Office formats.
        """
        handlers = {
            # Standard formats
            "application/pdf": self._handle_pdf,
            "text/plain": self._handle_text,
            "image/png": self._handle_image,
            "image/jpeg": self._handle_image,
            "text/csv": self._handle_text,

            # Office formats (these are what Google Docs export to)
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": self._handle_docx,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": self._handle_xlsx,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": self._handle_pptx,

            # Google native types (these should generally be exported before handling)
            "application/vnd.google-apps.document": self._handle_google_doc,
            "application/vnd.google-apps.spreadsheet": self._handle_google_sheet,
            "application/vnd.google-apps.presentation": self._handle_google_slide,
        }

        # If we don't have a specific handler, try to guess based on MIME type patterns
        if mime_type not in handlers:
            if "wordprocessing" in mime_type or "document" in mime_type:
                return self._handle_docx
            elif "spreadsheet" in mime_type or "sheet" in mime_type:
                return self._handle_xlsx
            elif "presentation" in mime_type or "slide" in mime_type:
                return self._handle_pptx
            elif "pdf" in mime_type:
                return self._handle_pdf
            elif "text" in mime_type:
                return self._handle_text
            elif "image" in mime_type:
                return self._handle_image

        return handlers.get(mime_type, self._handle_unknown)

    def _convert_structured_to_text(self, mime_type: str, structured_data: Dict[str, Any]) -> Optional[str]:
        """
        Converts structured data to a text representation for file types
        that do not have a direct text extraction. For example, XLSX files are
        processed into sheet-by-sheet text output.
        """
        if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            # Assume structured_data is a dict with sheet names as keys and lists of rows as values.
            text_lines = []
            for sheet_name, rows in structured_data.items():
                text_lines.append(f"Sheet: {sheet_name}")
                for row in rows:
                    # Convert each cell in the row to string and join with a tab.
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    text_lines.append(row_text)
                text_lines.append("")  # Blank line between sheets.
            return "\n".join(text_lines)
        # Extend this method for other MIME types as needed.
        return None

    # -------------------------------------------------------------------------
    # Extraction Handlers for Different MIME Types
    # -------------------------------------------------------------------------
    async def _handle_pdf(self, content: bytes) -> dict:
        """
        Extracts text and tables from a PDF file using pdfplumber.
        """
        text = ""
        structured = []
        try:
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                    tables = page.extract_tables()
                    if tables:
                        structured.append({"page": page.page_number, "tables": tables})
        except Exception as e:
            self.logger.error(f"Error extracting PDF content: {e}")
        return {
            "extracted_text": text.strip() if text else None,
            "structured_data": structured if structured else None
        }

    async def _handle_docx(self, content: bytes) -> dict:
        """
        Extracts text and tables from a DOCX file using python-docx.
        """
        text = ""
        tables = []
        try:
            doc = Document(io.BytesIO(content))
            paragraphs = [para.text for para in doc.paragraphs if para.text]
            text = "\n".join(paragraphs)
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
        except Exception as e:
            self.logger.error(f"Error extracting DOCX content: {e}")
        return {
            "extracted_text": text if text else None,
            "structured_data": {"tables": tables} if tables else None
        }

    async def _handle_xlsx(self, content: bytes) -> dict:
        """
        Extracts data from an XLSX file using openpyxl.
        Returns the data as structured data only; the conversion to text
        is performed later by _convert_structured_to_text.
        """
        sheets = {}
        try:
            wb = load_workbook(io.BytesIO(content), data_only=True)
            for sheet in wb.worksheets:
                sheet_data = []
                for row in sheet.iter_rows(values_only=True):
                    sheet_data.append(list(row))
                sheets[sheet.title] = sheet_data
        except Exception as e:
            self.logger.error(f"Error extracting XLSX content: {e}")
        return {
            "extracted_text": None,
            "structured_data": sheets if sheets else None
        }

    async def _handle_text(self, content: bytes) -> dict:
        """
        Decodes plain text file bytes to a string.
        """
        try:
            decoded_text = content.decode("utf-8", errors="replace")
        except Exception as e:
            self.logger.error(f"Error decoding text content: {e}")
            decoded_text = None
        return {
            "extracted_text": decoded_text,
            "structured_data": None
        }

    async def _handle_image(self, content: bytes) -> dict:
        """
        Returns a base64-encoded string of the image.
        (You can also add OCR extraction here if desired.)
        """
        try:
            encoded = base64.b64encode(content).decode("utf-8")
        except Exception as e:
            self.logger.error(f"Error encoding image content to base64: {e}")
            encoded = None
        return {
            "extracted_text": encoded,
            "structured_data": None
        }

    async def _handle_pptx(self, content: bytes) -> dict:
        """
        Extracts text from a PPTX file using python-pptx.
        """
        extracted_text = []
        try:
            prs = Presentation(io.BytesIO(content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        extracted_text.append(shape.text)
        except Exception as e:
            self.logger.error(f"Error extracting PPTX content: {e}")
        full_text = "\n".join(extracted_text)
        return {
            "extracted_text": full_text.strip() if full_text else None,
            "structured_data": None
        }

    async def _handle_google_doc(self, content: bytes) -> dict:
        """
        Handler for Google Docs content.
        First tries to process it as a DOCX (exported format),
        and if that fails, falls back to the warning.
        """
        try:
            # Try to handle it as a DOCX file first (the export format)
            docx_result = await self._handle_docx(content)
            if docx_result["extracted_text"] or docx_result["structured_data"]:
                return docx_result

            # If that didn't work, try as plain text
            text_result = await self._handle_text(content)
            if text_result["extracted_text"]:
                return text_result

            # If we still have no content, log warning
            self.logger.warning("Google Docs extraction attempted but no content extracted")
            return {
                "extracted_text": None,
                "structured_data": None
            }
        except Exception as e:
            self.logger.error(f"Error handling Google Doc content: {e}")
            return {
                "extracted_text": None,
                "structured_data": None
            }

    # Replace the Google Sheet handler
    async def _handle_google_sheet(self, content: bytes) -> dict:
        """
        Handler for Google Sheets content.
        First tries to process it as XLSX (exported format),
        and if that fails, falls back to the warning.
        """
        try:
            # Try to handle it as an XLSX file first (the export format)
            xlsx_result = await self._handle_xlsx(content)
            if xlsx_result["structured_data"]:
                return xlsx_result

            # If that didn't work, try as text/csv
            text_result = await self._handle_text(content)
            if text_result["extracted_text"]:
                return text_result

            # If we still have no content, log warning
            self.logger.warning("Google Sheets extraction attempted but no content extracted")
            return {
                "extracted_text": None,
                "structured_data": None
            }
        except Exception as e:
            self.logger.error(f"Error handling Google Sheet content: {e}")
            return {
                "extracted_text": None,
                "structured_data": None
            }

    # Replace the Google Slide handler
    async def _handle_google_slide(self, content: bytes) -> dict:
        """
        Handler for Google Slides content.
        First tries to process it as PPTX (exported format),
        and if that fails, falls back to the warning.
        """
        try:
            # Try to handle it as a PPTX file first (the export format)
            pptx_result = await self._handle_pptx(content)
            if pptx_result["extracted_text"]:
                return pptx_result

            # If we still have no content, log warning
            self.logger.warning("Google Slides extraction attempted but no content extracted")
            return {
                "extracted_text": None,
                "structured_data": None
            }
        except Exception as e:
            self.logger.error(f"Error handling Google Slide content: {e}")
            return {
                "extracted_text": None,
                "structured_data": None
            }

    # 3. Update the extract method to be more robust with MIME type detection
    async def extract(self, file_id: str, content: bytes, mime_type: str = "") -> ExtractedContent:
        """
        Detects the MIME type (if not provided) and calls the appropriate handler
        to extract content from the file. If no direct text is returned but structured
        data is available (for example, in XLSX files), this method converts the
        structured data to a text representation.
        """
        if not content:
            self.logger.warning(f"No content provided for file {file_id}")
            return ExtractedContent(
                file_id=file_id,
                mime_type=mime_type or "unknown",
                extracted_text=None,
                structured_data=None,
                raw_content=None
            )

        # Determine MIME type if not provided
        detected_mime_type = mime_type
        if not detected_mime_type:
            try:
                detected_mime_type = self.mime_detector.from_buffer(content)
            except Exception as e:
                self.logger.error(f"Error detecting MIME type for file {file_id}: {e}")
                detected_mime_type = "application/octet-stream"  # Generic binary

        # Map Google MIME types to their exported equivalents if we still have a Google type
        google_mime_map = {
            "application/vnd.google-apps.document": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.google-apps.spreadsheet": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.google-apps.presentation": "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        }

        # Use the mapped MIME type if applicable
        processing_mime_type = detected_mime_type
        if detected_mime_type in google_mime_map:
            self.logger.info(
                f"Using exported format for Google file: {detected_mime_type} -> {google_mime_map[detected_mime_type]}")
            processing_mime_type = google_mime_map[detected_mime_type]

        # Retrieve the handler based on the MIME type
        handler = self._get_handler(processing_mime_type)

        try:
            result = await handler(content)
            extracted_text = result.get("extracted_text")
            structured_data = result.get("structured_data")

            # If no text was directly extracted but structured data exists,
            # convert that structured data into a text representation
            if not extracted_text and structured_data:
                converted_text = self._convert_structured_to_text(processing_mime_type, structured_data)
                if converted_text:
                    extracted_text = converted_text

            return ExtractedContent(
                file_id=file_id,
                mime_type=detected_mime_type,  # Return the original detected MIME type
                extracted_text=extracted_text,
                structured_data=structured_data,
                raw_content=content
            )
        except Exception as e:
            self.logger.error(f"Error extracting content for file {file_id}: {e}")
            return ExtractedContent(
                file_id=file_id,
                mime_type=detected_mime_type,
                extracted_text=None,
                structured_data=None,
                raw_content=content
            )




    async def _handle_unknown(self, content: bytes) -> dict:
        """
        Fallback handler for unknown or unsupported file types.
        """
        self.logger.warning("Unknown MIME type. No extraction performed.")
        return {
            "extracted_text": None,
            "structured_data": None
        }


# -----------------------------------------------------------------------------
# Example Usage
# -----------------------------------------------------------------------------
if __name__ == "__main__":
    import sys

    logging.basicConfig(stream=sys.stdout, level=logging.DEBUG)


    async def main():
        # Replace with your own file_id and file path.
        file_id = "example_file_id"
        file_path = "AndhraPradesh_16082024.xlsx"  # Change to an appropriate test file.
        try:
            with open(file_path, "rb") as f:
                file_bytes = f.read()
        except Exception as e:
            print(f"Failed to open file: {e}")
            return
        logger = Logger()
        extractor = ContentExtractor(logger=logger)
        extracted = await extractor.extract(file_id, file_bytes)
        print("MIME Type:", extracted.mime_type)
        print("Extracted Text:\n", extracted.extracted_text)
        print("Structured Data:\n", extracted.structured_data)


    asyncio.run(main())
